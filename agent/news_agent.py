from langchain.chains import TransformChain, SequentialChain
from langchain.agents import AgentExecutor, Tool
from langchain_core.runnables import RunnableParallel
import os


class NewsPPTAgent:
    def __init__(self):
        # 初始化子链
        self.process_chain = self._create_process_chain()  # 创建新闻预处理链
        self.audio_chain = self._create_audio_chain()      # 创建音频生成链
        self.image_chain = self._create_image_chain()      # 创建图像生成链
        self.ppt_chain = self._create_ppt_chain()          # 创建PPT组装链

        # 构建完整流程
        self.full_chain = RunnableParallel(
            processed_news=self.process_chain,  # 新闻预处理结果
            audio_files=self.audio_chain,       # 音频文件路径
            images=self.image_chain             # 图像文件路径
        ) | self.ppt_chain  # 将所有链按顺序连接

    def _create_process_chain(self):
        def news_processor(inputs):
            # 新闻预处理逻辑
            # 这里可以添加具体的新闻预处理代码，例如提取关键词和生成摘要
            return {"keywords": ["关键词1", "关键词2"], "summary": "新闻摘要"}

        return TransformChain(
            input_variables=["news"],  # 输入变量
            output_variables=["processed_news"],  # 输出变量
            transform=news_processor  # 处理函数
        )

    def _create_audio_chain(self):
        def audio_generator(inputs):
            # 音频生成逻辑
            # 这里可以添加具体的音频生成代码，例如调用语音合成API
            return {"audio_path": "audio.mp3"}

        return TransformChain(
            input_variables=["processed_news"],  # 输入变量
            output_variables=["audio_files"],  # 输出变量
            transform=audio_generator  # 处理函数
        )

    def _create_image_chain(self):
        def image_generator(inputs):
            # 宫崎骏风格图片生成逻辑
            # 这里可以添加具体的图像生成代码，例如调用图像生成API
            return {"image_path": "image.png"}

        return TransformChain(
            input_variables=["processed_news"],  # 输入变量
            output_variables=["images"],  # 输出变量
            transform=image_generator  # 处理函数
        )

    def _create_ppt_chain(self):
        def ppt_assembler(inputs):
            # PPT组装逻辑
            # 这里可以添加具体的PPT组装代码，例如调用PPT生成库
            return {"ppt_path": "news.pptx"}

        return TransformChain(
            input_variables=["processed_news", "audio_files", "images"],  # 输入变量
            output_variables=["ppt_path"],  # 输出变量
            transform=ppt_assembler  # 处理函数
        )

    def run(self, news_list):
        # 批量处理新闻列表
        return self.full_chain.batch([{"news": news} for news in news_list])


from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.llms import OpenAI  # 使用 langchain.llms.OpenAI 替代 langchain_openai.OpenAI


def create_miyazaki_image_chain():
    prompt = PromptTemplate(
        input_variables=["keywords", "summary"],  # 输入变量
        template="""生成宫崎骏风格插画提示词：
        要求：吉卜力动画风格、柔和水彩质感、奇幻场景
        新闻关键词：{keywords}
        新闻摘要：{summary}
        输出格式：英文提示词，包含'miyazaki hayao style'"""
    )

    return LLMChain(
        llm=OpenAI(temperature=0.7),  # 使用 OpenAI 模型，温度设置为 0.7
        prompt=prompt,  # 提示模板
        output_key="image_prompt"  # 输出键
    )


def generate_image(image_prompt):
    from openai import OpenAI
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))  # 使用环境变量获取 API 密钥

    response = client.images.generate(
        model="dall-e-3",  # 使用 DALL-E 3 模型
        prompt=image_prompt + ", Studio Ghibli style, watercolor texture, fantasy elements",  # 提示词
        size="1024x1024",  # 图像大小
        quality="standard",  # 图像质量
        n=1  # 生成的图像数量
    )
    return response.data[0].url  # 返回生成的图像 URL


from langchain.tools import tool
from openai import OpenAI


@tool
def generate_news_audio(text: str, output_dir: str = "audio"):
    """生成新闻语音并保存为MP3"""
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))  # 使用环境变量获取 API 密钥

    os.makedirs(output_dir, exist_ok=True)  # 创建输出目录
    filename = f"{hash(text)}.mp3"  # 生成文件名
    path = os.path.join(output_dir, filename)  # 文件路径

    response = client.audio.speech.create(
        model="tts-1",  # 使用 TTS-1 模型
        voice="nova",  # 适合新闻播报的声音
        input=text  # 输入文本
    )
    response.stream_to_file(path)  # 将音频流保存为文件
    return path


from pptx import Presentation
from pptx.util import Inches
import requests
from PIL import Image
import io


def create_ppt_page(slide, title, image_path, audio_path):
    # 添加标题
    title_shape = slide.shapes.title
    title_shape.text = title

    # 添加图片
    img = Image.open(requests.get(image_path, stream=True).raw)  # 下载图像
    img_io = io.BytesIO()
    img.save(img_io, format='PNG')  # 将图像保存为 PNG 格式

    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(img_io, left, top, width=Inches(6))  # 添加图片到幻灯片

    # 添加音频（Windows系统可用）
    if os.name == 'nt':
        slide.shapes.add_movie(
            audio_path,
            left, Inches(5),
            width=Inches(1),
            height=Inches(1))  # 添加音频到幻灯片


def generate_ppt(pages_data, output_path="news_report.pptx"):
    prs = Presentation()  # 创建新的 PPT 演示文稿
    for data in pages_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 添加新幻灯片
        create_ppt_page(
            slide,
            title=data['title'],  # 幻灯片标题
            image_path=data['image_url'],  # 图像路径
            audio_path=data['audio_path']  # 音频路径
        )
    prs.save(output_path)  # 保存 PPT 文件
