from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# === 配置参数 ===
TEMPLATE_PATH = "./ppt/a.pptx"  # 模板路径
OUTPUT_PATH = "news_presentation_left_image_right_text.pptx"
NEWS_CONTENTS = [
    {
        "title": "新闻一：量子计算机突破",
        "content": "科学家近日宣布成功研发出新型量子计算机，运算速度较传统计算机提升千倍，有望彻底改变人工智能领域。",
        "image_path": "images/quantum.png",  # 替换为实际图片路径
        "audio_path": "audios/1.mp3"  # 替换为实际音频路径
    },
    {
        "title": "新闻二：环保新协议",
        "content": "全球多个国家签署《巴黎协定》补充协议，承诺2030年前将碳排放量减少50%，并推动可再生能源发展。",
        "image_path": "images/environment.png",
        "audio_path": "audios/2.mp3"
    }
]

def add_news_slide(prs, title, content, image_path, audio_path):
    layout = prs.slide_layouts[6]  # 空白版式
    slide = prs.slides.add_slide(layout)

    # 图片区域（左侧）
    img_left = Inches(0.5)
    img_top = Inches(1.5)
    img_width = Inches(4.5)
    img_height = Inches(3)

    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            left=img_left,
            top=img_top,
            width=img_width,
            height=img_height
        )

    # 文字区域（右侧）
    text_left = Inches(5.5)
    text_top = Inches(1.5)
    text_width = Inches(6)
    text_height = Inches(3)

    text_box = slide.shapes.add_textbox(
        left=text_left,
        top=text_top,
        width=text_width,
        height=text_height
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    p_title = tf.add_paragraph()
    p_title.text = f"【{title}】"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.alignment = PP_ALIGN.LEFT

    p_content = tf.add_paragraph()
    p_content.text = content
    p_content.font.size = Pt(22)
    p_content.alignment = PP_ALIGN.LEFT

    # 添加音频占位符（隐藏的矩形）
    placeholder_shape = slide.shapes.add_shape(
        autoshape_type_id=MSO_SHAPE.RECTANGLE,  # 添加此参数
        left=Inches(-1),  # 隐藏形状
        top=Inches(-1),
        width=Inches(0.1),
        height=Inches(0.1)
    )
    placeholder_shape.name = f"AudioPlaceholder_{audio_path}"  # 设置形状名称作为音频标识

    return slide


def create_presentation():
    prs = Presentation(TEMPLATE_PATH)

    for news in NEWS_CONTENTS:
        add_news_slide(
            prs,
            news["title"],
            news["content"],
            news.get("image_path", ""),
            news.get("audio_path", "")
        )

    prs.save(OUTPUT_PATH)
    print(f"已生成 {len(NEWS_CONTENTS)} 页图文混合PPT，保存至: {OUTPUT_PATH}")




if __name__ == "__main__":
    create_presentation()